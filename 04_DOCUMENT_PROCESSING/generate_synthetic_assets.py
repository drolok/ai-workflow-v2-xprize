from __future__ import annotations

from pathlib import Path

from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(r"C:\AI_WORKFLOW_V2\04_DOCUMENT_PROCESSING")
INBOX = ROOT / "00_Inbox"


def ensure_inbox() -> None:
    INBOX.mkdir(parents=True, exist_ok=True)


def pick_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path(r"C:\Windows\Fonts\arial.ttf"),
        Path(r"C:\Windows\Fonts\segoeui.ttf"),
        Path(r"C:\Windows\Fonts\calibri.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def build_html_source() -> Path:
    html_path = INBOX / "phase4_pdf_source.html"
    html = """<!doctype html>
<html lang="en">
  <head>
    <meta charset="utf-8" />
    <title>Phase 4 PDF Synthetic Test</title>
    <style>
      body {
        font-family: "Segoe UI", Arial, sans-serif;
        margin: 48px;
        color: #1b1f24;
        line-height: 1.5;
      }
      h1 {
        margin-bottom: 0.2rem;
      }
      .note {
        background: #f4f7fb;
        border-left: 6px solid #1f6feb;
        padding: 12px 16px;
        margin: 20px 0;
      }
      table {
        border-collapse: collapse;
        width: 100%;
        margin-top: 16px;
      }
      th, td {
        border: 1px solid #b8c2cc;
        padding: 8px 10px;
        text-align: left;
      }
      th {
        background: #e9eef5;
      }
      code {
        background: #f6f8fa;
        padding: 2px 5px;
      }
    </style>
  </head>
  <body>
    <h1>Phase 4 PDF Synthetic Test</h1>
    <p>This PDF exists only to validate the AI_WORKFLOW_V2 document pipeline.</p>
    <div class="note">
      <strong>Project code:</strong> <code>LANTERN-42</code>
    </div>
    <h2>Key facts</h2>
    <ul>
      <li>The storage mode is local-first.</li>
      <li>The next planned framework phase after this one is automation.</li>
      <li>This file was generated on 2026-07-18 for synthetic testing only.</li>
    </ul>
    <h2>Status table</h2>
    <table>
      <thead>
        <tr>
          <th>Component</th>
          <th>Status</th>
        </tr>
      </thead>
      <tbody>
        <tr>
          <td>Docling pipeline</td>
          <td>Validation target</td>
        </tr>
        <tr>
          <td>Markdown output</td>
          <td>Required</td>
        </tr>
      </tbody>
    </table>
  </body>
</html>
"""
    html_path.write_text(html, encoding="utf-8")
    return html_path


def build_png_source() -> Path:
    png_path = INBOX / "phase4_ocr_image.png"
    image = Image.new("RGB", (1600, 520), "white")
    draw = ImageDraw.Draw(image)
    font_large = pick_font(78)
    font_medium = pick_font(52)

    draw.text((60, 70), "OCR TOKEN", fill="black", font=font_large)
    draw.text((60, 200), "HELIOS-17", fill="black", font=font_large)
    draw.text((60, 340), "phase 4 synthetic image", fill="black", font=font_medium)

    image.save(png_path)
    return png_path


def build_docx_source() -> Path:
    docx_path = INBOX / "phase4_docx_test.docx"
    doc = Document()
    doc.add_heading("Phase 4 DOCX Synthetic Test", level=1)
    doc.add_paragraph("This DOCX is synthetic and exists only for Docling validation.")
    doc.add_paragraph("Key token: MONDRIAN-88")
    doc.add_paragraph("The intended destination format is Markdown.")
    doc.add_paragraph("The framework remains local-first and localhost-only.")
    doc.save(docx_path)
    return docx_path


def build_pptx_source() -> Path:
    pptx_path = INBOX / "phase4_pptx_test.pptx"
    deck = Presentation()

    title_layout = deck.slide_layouts[0]
    slide = deck.slides.add_slide(title_layout)
    slide.shapes.title.text = "Phase 4 PPTX Synthetic Test"
    slide.placeholders[1].text = "Token: CONSTELLATION-9"

    bullets_layout = deck.slide_layouts[1]
    slide = deck.slides.add_slide(bullets_layout)
    slide.shapes.title.text = "Pipeline facts"
    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()
    first = text_frame.paragraphs[0]
    first.text = "This presentation is synthetic."
    first.font.size = Pt(22)

    for bullet in [
        "Docling should convert PPTX into readable Markdown.",
        "No real user data is included.",
        "The next major phase is automation.",
    ]:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)

    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Reference"
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(8.0), Inches(2.0))
    frame = textbox.text_frame
    frame.text = "Validation phrase: ORBIT-CHECK-5"
    frame.paragraphs[0].font.size = Pt(28)

    deck.save(pptx_path)
    return pptx_path


def main() -> None:
    ensure_inbox()
    created = [
        build_html_source(),
        build_png_source(),
        build_docx_source(),
        build_pptx_source(),
    ]
    for path in created:
        print(path)


if __name__ == "__main__":
    main()
